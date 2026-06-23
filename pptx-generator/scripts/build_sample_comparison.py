#!/usr/bin/env python3
"""Generate the sample_comparison.pptx file from editable Python source.

Regeneration:
    python -m pip install python-pptx
    python pptx-generator/scripts/build_sample_comparison.py

The generated PPTX is written to pptx-generator/output/sample_comparison.pptx.
The SVG preview is written to pptx-generator/preview/sample_comparison.svg.
The PPTX is intentionally not committed because GitHub PR file preview does not
support binary PPTX files. Commit the script and SVG preview instead.
"""
from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency guidance for CLI users
    raise SystemExit(
        "python-pptx is required to regenerate the PPTX.\n"
        "Install it with: python -m pip install python-pptx"
    ) from exc

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "output" / "sample_comparison.pptx"
PREVIEW_PATH = ROOT / "preview" / "sample_comparison.svg"

TEAL = RGBColor(0x00, 0xAF, 0xB2)
PALE_TEAL = RGBColor(0xEA, 0xFB, 0xFB)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xC8, 0xEF, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def add_card(slide, left, top, width, height, title, body, number=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1.2)

    if number is not None:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), top + Inches(0.18), Inches(0.38), Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.color.rgb = TEAL
        text_frame = badge.text_frame
        text_frame.text = str(number)
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        title_left = left + Inches(0.68)
    else:
        title_left = left + Inches(0.25)

    add_textbox(slide, title_left, top + Inches(0.17), width - Inches(0.45), Inches(0.28), title, 13, True, TEAL)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.58), width - Inches(0.5), height - Inches(0.72), body, 10, False, TEXT)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE

    add_textbox(slide, Inches(0.65), Inches(0.45), Inches(8.5), Inches(0.48), "生成AI研修：改善前後サンプル", 25, True, TEXT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.08), Inches(5.8), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.color.rgb = TEAL

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.35), Inches(0.42), Inches(2.25), Inches(0.42))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = PALE_TEAL
    ribbon.line.color.rgb = BORDER
    add_textbox(slide, Inches(10.58), Inches(0.51), Inches(1.8), Inches(0.2), "PPTX DESIGN", 9, True, TEAL)

    add_textbox(slide, Inches(0.75), Inches(1.48), Inches(5.4), Inches(0.3), "Before：情報量が多く視線が迷う", 15, True, MUTED)
    add_textbox(slide, Inches(7.15), Inches(1.48), Inches(5.4), Inches(0.3), "After：要点をカード化して理解しやすく", 15, True, TEAL)

    before = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.95), Inches(5.85), Inches(4.85))
    before.fill.solid()
    before.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    before.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    add_textbox(slide, Inches(1.0), Inches(2.28), Inches(5.1), Inches(0.35), "AI活用で業務を効率化するポイント", 15, True, TEXT)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.85),
        Inches(5.0),
        Inches(2.4),
        "・作業の目的を明確にする\n・入力情報を整理する\n・AIの回答を確認する\n・社内ルールに合わせて活用する\n・小さく試して改善する",
        13,
        False,
        TEXT,
    )
    add_textbox(slide, Inches(1.0), Inches(5.75), Inches(4.9), Inches(0.46), "課題：箇条書きが並び、重要度の差が見えにくい", 11, False, MUTED)

    after_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.95), Inches(5.85), Inches(4.85))
    after_bg.fill.solid()
    after_bg.fill.fore_color.rgb = PALE_TEAL
    after_bg.line.color.rgb = BORDER
    add_card(slide, Inches(7.2), Inches(2.35), Inches(2.45), Inches(1.52), "目的を決める", "何を短縮したいかを先に決めます。", 1)
    add_card(slide, Inches(9.9), Inches(2.35), Inches(2.45), Inches(1.52), "情報を渡す", "背景・条件・期待する形式を伝えます。", 2)
    add_card(slide, Inches(7.2), Inches(4.25), Inches(2.45), Inches(1.52), "回答を確認", "事実・表現・社内ルールを見直します。", 3)
    add_card(slide, Inches(9.9), Inches(4.25), Inches(2.45), Inches(1.52), "小さく改善", "試行結果を次の指示に反映します。", 4)
    add_textbox(slide, Inches(7.25), Inches(6.15), Inches(5.0), Inches(0.34), "効果：短い見出しと配置で、行動手順がすぐ分かる", 11, False, TEAL)

    return prs


def build_svg_preview() -> str:
    """Return an SVG preview that mirrors the generated sample slide."""
    return '''<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720" role="img" aria-labelledby="title desc">
  <title id="title">生成AI研修：改善前後サンプル</title>
  <desc id="desc">PPTX生成スクリプトで作成される比較スライドのSVGプレビュー。</desc>
  <rect width="1280" height="720" fill="#ffffff"/>
  <text x="62" y="78" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="36" font-weight="700" fill="#1f2937">生成AI研修：改善前後サンプル</text>
  <rect x="62" y="104" width="557" height="5" rx="2.5" fill="#00afb2"/>
  <rect x="994" y="40" width="216" height="40" rx="14" fill="#eafbfb" stroke="#c8efef"/>
  <text x="1054" y="66" font-family="Arial, sans-serif" font-size="14" font-weight="700" fill="#00afb2">PPTX DESIGN</text>

  <text x="72" y="171" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="22" font-weight="700" fill="#64748b">Before：情報量が多く視線が迷う</text>
  <rect x="62" y="187" width="562" height="466" rx="22" fill="#f8fafc" stroke="#cbd5e1"/>
  <text x="96" y="247" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="23" font-weight="700" fill="#1f2937">AI活用で業務を効率化するポイント</text>
  <g font-family="Arial, 'Noto Sans JP', sans-serif" font-size="23" fill="#1f2937">
    <text x="96" y="309">・作業の目的を明確にする</text>
    <text x="96" y="353">・入力情報を整理する</text>
    <text x="96" y="397">・AIの回答を確認する</text>
    <text x="96" y="441">・社内ルールに合わせて活用する</text>
    <text x="96" y="485">・小さく試して改善する</text>
  </g>
  <text x="96" y="596" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="18" fill="#64748b">課題：箇条書きが並び、重要度の差が見えにくい</text>

  <text x="686" y="171" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="22" font-weight="700" fill="#00afb2">After：要点をカード化して理解しやすく</text>
  <rect x="658" y="187" width="562" height="466" rx="22" fill="#eafbfb" stroke="#c8efef"/>
  <g font-family="Arial, 'Noto Sans JP', sans-serif">
    <g transform="translate(691 226)">
      <rect width="235" height="146" rx="18" fill="#fff" stroke="#c8efef"/>
      <circle cx="36" cy="36" r="18" fill="#00afb2"/><text x="31" y="42" font-size="16" font-weight="700" fill="#fff">1</text>
      <text x="66" y="42" font-size="19" font-weight="700" fill="#00afb2">目的を決める</text>
      <text x="24" y="89" font-size="16" fill="#1f2937">何を短縮したいかを</text><text x="24" y="115" font-size="16" fill="#1f2937">先に決めます。</text>
    </g>
    <g transform="translate(950 226)">
      <rect width="235" height="146" rx="18" fill="#fff" stroke="#c8efef"/>
      <circle cx="36" cy="36" r="18" fill="#00afb2"/><text x="31" y="42" font-size="16" font-weight="700" fill="#fff">2</text>
      <text x="66" y="42" font-size="19" font-weight="700" fill="#00afb2">情報を渡す</text>
      <text x="24" y="89" font-size="16" fill="#1f2937">背景・条件・期待する</text><text x="24" y="115" font-size="16" fill="#1f2937">形式を伝えます。</text>
    </g>
    <g transform="translate(691 408)">
      <rect width="235" height="146" rx="18" fill="#fff" stroke="#c8efef"/>
      <circle cx="36" cy="36" r="18" fill="#00afb2"/><text x="31" y="42" font-size="16" font-weight="700" fill="#fff">3</text>
      <text x="66" y="42" font-size="19" font-weight="700" fill="#00afb2">回答を確認</text>
      <text x="24" y="89" font-size="16" fill="#1f2937">事実・表現・社内ルールを</text><text x="24" y="115" font-size="16" fill="#1f2937">見直します。</text>
    </g>
    <g transform="translate(950 408)">
      <rect width="235" height="146" rx="18" fill="#fff" stroke="#c8efef"/>
      <circle cx="36" cy="36" r="18" fill="#00afb2"/><text x="31" y="42" font-size="16" font-weight="700" fill="#fff">4</text>
      <text x="66" y="42" font-size="19" font-weight="700" fill="#00afb2">小さく改善</text>
      <text x="24" y="89" font-size="16" fill="#1f2937">試行結果を次の指示に</text><text x="24" y="115" font-size="16" fill="#1f2937">反映します。</text>
    </g>
  </g>
  <text x="696" y="617" font-family="Arial, 'Noto Sans JP', sans-serif" font-size="18" fill="#00afb2">効果：短い見出しと配置で、行動手順がすぐ分かる</text>
</svg>
'''


def save_svg_preview() -> None:
    PREVIEW_PATH.parent.mkdir(parents=True, exist_ok=True)
    PREVIEW_PATH.write_text(build_svg_preview(), encoding="utf-8")


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    build_deck().save(OUTPUT_PATH)
    save_svg_preview()
    print(f"Generated {OUTPUT_PATH}")
    print(f"Generated {PREVIEW_PATH}")


if __name__ == "__main__":
    main()
